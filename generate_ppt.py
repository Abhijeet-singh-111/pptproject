from flask import Flask, request, send_file, jsonify
from flask_cors import CORS # Import CORS
from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.util import Inches, Pt
import re
import html
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io # To handle file in memory
import sys # Import sys for printing debug messages to stderr

# --- Helper Functions ---

# Function to extract formatted text runs from a node, respecting display:none and inline styles
# This function checks styles on *each* tag as it traverses
def extract_formatted_text_runs(element, default_font_size_pt, default_color_rgb, default_is_bold=False):
    runs_data = []
    if element is None:
        return runs_data

    # Determine the current inherited styles based on the element itself
    element_style = element.get('style', '')
    current_font_size_pt = default_font_size_pt
    current_color_rgb = default_color_rgb
    current_is_bold = default_is_bold

    parsed_size_px = parse_font_size_px(element_style)
    if parsed_size_px:
        current_font_size_pt = Pt(parsed_size_px * 0.75) # Override default if style is present

    parsed_color = parse_color(element_style)
    if parsed_color:
        current_color_rgb = parsed_color # Override default if style is present

    element_is_bold_override = parse_font_weight(element_style)
    if 'font-weight:' in element_style: # Check if font-weight is explicitly set
         current_is_bold = element_is_bold_override # Override default if style is present


    for content in element.contents:
        if isinstance(content, NavigableString): # This is a text node
            text = str(content)
            # Strip leading/trailing whitespace from the text node content
            cleaned_text = text.strip()
            # Replace internal sequences of whitespace (including newlines) with a single space
            cleaned_text = re.sub(r'\s+', ' ', cleaned_text)

            if cleaned_text: # Only add as a part if there's non-whitespace content
                 # Use the current styles derived from the element's own style or inherited
                 runs_data.append({
                     'text': html.unescape(cleaned_text),
                     'is_superscript': False,
                     'font_size': current_font_size_pt,
                     'color': current_color_rgb,
                     'is_bold': current_is_bold
                 })

        elif isinstance(content, Tag): # This is an element node
            style = content.get('style', '')
            # Check if the element itself has display:none
            if 'display:none' in style.replace(" ", ""):
                 continue # Skip this element and its children if display is none

            # Handle specific tags that might add formatting (like <sup>, <br>)
            if content.name == 'sup':
                # For superscript, apply superscript flag and recurse
                sup_content_runs = extract_formatted_text_runs(
                    content,
                    current_font_size_pt, # Pass down the current size
                    current_color_rgb, # Pass down the current color
                    current_is_bold # Pass down the current bold
                )
                for run_data in sup_content_runs:
                    run_data['is_superscript'] = True # Mark as superscript
                    # Optional: Reduce size for superscript content here if not handled by style
                    # if 'font_size' not in run_data and run_data['font_size'] is not None:
                    #      run_data['font_size'] = Pt(run_data['font_size'].pt * 0.7)
                    runs_data.append(run_data)
            elif content.name == 'br':
                 # Handle line breaks as a separate run
                 runs_data.append({'text': '\n', 'is_superscript': False, 'font_size': current_font_size_pt, 'color': current_color_rgb, 'is_bold': current_is_bold})
            else:
                # For other tags, recursively process children.
                # The styles *within* the recursive call will be determined by the child's style
                # overriding the styles passed down from the parent (current_...).
                runs_data.extend(extract_formatted_text_runs(
                    content,
                    current_font_size_pt, # Pass current inherited styles as defaults for children
                    current_color_rgb,
                    current_is_bold
                ))

    return runs_data


# Function to parse pixel font size from style string
def parse_font_size_px(style_string):
    font_size_match = re.search(r'font-size:\s*(\d+)px', style_string)
    if font_size_match:
        return int(font_size_match.group(1))
    return None

# Function to parse color from style string (basic)
def parse_color(style_string):
    color_match = re.search(r'color:\s*([^;]+)', style_string)
    if color_match:
        color_name_or_hex = color_match.group(1).strip()
        if color_name_or_hex.lower() == 'red':
             return RGBColor(0xFF, 0x00, 0x00)
        elif color_name_or_hex.startswith('#') and len(color_name_or_hex) == 7:
             try:
                 rgb_tuple = bytes.fromhex(color_name_or_hex[1:])
                 return RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])
             except ValueError:
                 pass # Invalid hex
        # Add more color name mappings if needed
    return None # Default or if parsing failed

# Function to parse font weight from style string
def parse_font_weight(style_string):
    if 'font-weight:600' in style_string or 'font-weight:bold' in style_string:
        return True
    return False

# Function to parse translate(x, y) from transform style (assuming px)
def parse_translate_px(style_string):
    translate_match = re.search(r'transform:\s*translate\(\s*(-?\d+\.?\d*)px\s*,\s*(-?\d+\.?\d*)px\s*\)', style_string)
    if translate_match:
        try:
            translate_x = float(translate_match.group(1))
            translate_y = float(translate_match.group(2))
            return translate_x, translate_y
        except ValueError:
            pass # Could not convert to float
    return 0, 0 # Return 0,0 if not found or parsing failed

# --- Constants for Pixel to Inch Conversion (assuming 96 DPI) ---
PX_TO_INCHES = 1.0 / 96.0

app = Flask(__name__)
CORS(app) # Apply CORS to the entire Flask application

@app.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    # Get HTML content from the request body (assuming JSON format)
    data = request.get_json()
    if not data or 'html_content' not in data:
        sys.stderr.write("Error: No HTML content provided in request.\n")
        return jsonify({"error": "No HTML content provided"}), 400

    html_string = data['html_content']

    try:
        # --- PPTX Generation Logic ---
        sys.stderr.write("Starting PPTX generation logic via Flask...\n") # Debug print to stderr
        prs = Presentation()
        # Using blank layout. Adjust if you need a different base layout from your template.
        blank_slide_layout = prs.slide_layouts[5]

        sys.stderr.write("Parsing HTML with BeautifulSoup...\n") # Debug print to stderr
        soup = BeautifulSoup(html_string, 'lxml')

        question_areas = soup.find_all('div', class_='QuestionArea')

        if not question_areas:
            sys.stderr.write("No 'QuestionArea' divs found in the HTML.\n") # Debug print to stderr
            # Return an error response if no content is found to process
            return jsonify({"error": "No question content found in HTML"}), 400
        else:
            sys.stderr.write(f"Found {len(question_areas)} 'QuestionArea' divs.\n") # Debug print to stderr


        for i, question_area in enumerate(question_areas):
            sys.stderr.write(f"Processing Question Area {i+1}...\n") # Debug print to stderr

            slide = prs.slides.add_slide(blank_slide_layout)

            # --- Remove the "Click to add title" placeholder if it exists ---
            for shape in slide.shapes:
                 if shape.is_placeholder:
                      try:
                           if shape.placeholder_format.idx == 0: # Index 0 is often the title placeholder
                                tf = shape.text_frame
                                tf.clear()
                                sys.stderr.write(f"  Slide {i+1}: Removed placeholder with index 0 (likely title).\n") # Debug print to stderr
                                break
                      except AttributeError:
                           pass

            # --- Extract and Add Question Text ---
            # Find the main question div (with red color and font size)
            question_main_div = question_area.find('div', style=lambda value: value and 'color: red' in value and 'font-size:30px' in value)

            if question_main_div:
                # Find the alterable parent div within the main question div (this is the one with transform)
                question_translated_div = question_main_div.find('div', class_='alterable parent')

                if question_translated_div:
                     question_style = question_translated_div.get('style', '')
                     translate_x_px, translate_y_px = parse_translate_px(question_style)

                     # Get default styles for the question block from the main question div
                     question_default_size_pt = Pt(20) # Default if not parsed
                     question_default_color_rgb = RGBColor(0,0,0) # Default black
                     question_default_is_bold = False

                     if question_main_div:
                           qs_style = question_main_div.get('style', '')
                           qs_parsed_size_px = parse_font_size_px(qs_style)
                           if qs_parsed_size_px:
                                question_default_size_pt = Pt(qs_parsed_size_px * 0.75)
                           qs_parsed_color = parse_color(qs_style)
                           if qs_parsed_color:
                                question_default_color_rgb = qs_parsed_color
                           question_default_is_bold = parse_font_weight(qs_style)


                     # Calculate final position with translate offset
                     base_left = Inches(0.5)
                     base_top = Inches(0.5)
                     width = Inches(9) # Adjust width as needed
                     height = Inches(1.5) # Adjust height as needed (can be adjusted after adding text)

                     final_left_inches = base_left.inches + (translate_x_px * PX_TO_INCHES)
                     final_top_inches = base_top.inches + (translate_y_px * PX_TO_INCHES)


                     # Extract formatted text runs from the translated question div, using question defaults
                     question_text_runs = extract_formatted_text_runs(
                          question_translated_div,
                          question_default_size_pt,
                          question_default_color_rgb,
                          question_default_is_bold
                     )

                     if question_text_runs:
                          sys.stderr.write(f"  Slide {i+1}: Adding Question shape.\n") # Debug print to stderr
                          txBox = slide.shapes.add_textbox(Inches(final_left_inches), Inches(final_top_inches), width, height)
                          tf = txBox.text_frame
                          tf.clear()
                          p = tf.add_paragraph()

                          for run_data in question_text_runs:
                               run = p.add_run()
                               run.text = run_data['text']
                               font = run.font
                               # Apply formatting from the run_data (extracted from inline styles or block defaults)
                               font.size = run_data['font_size'] # Use the size determined by extract_formatted_text_runs
                               font.bold = run_data['is_bold'] # Use the bold determined by extract_formatted_text_runs
                               if run_data.get('color'):
                                    font.color.rgb = run_data['color'] # Use the color determined by extract_formatted_text_runs
                               if run_data.get('is_superscript'):
                                    font.superscript = True
                                    # Optional: Reduce size for superscript if needed and not handled by inline style
                                    # if font.size is not None:
                                    #      font.size = Pt(font.size.pt * 0.7)


                else:
                     sys.stderr.write(f"  Slide {i+1}: Warning: Could not find the question's alterable parent div (for transform).\n")


            # --- Extract and Add Options ---
            options_container_div = question_area.find('div', style=lambda value: value and 'width:90%' in value and 'font-size:30px' in value)

            if options_container_div:
                 # Find all option divs (each with display:inline-flex and alterable parent)
                 option_elements = options_container_div.find_all('div', style=lambda value: value and 'display:inline-flex' in value and 'alterable parent' in value.split()) # Ensure it's also alterable parent

                 # Get default styles for option blocks from the options container div
                 options_default_size_pt = Pt(20) # Default if not parsed
                 options_default_color_rgb = RGBColor(0,0,0) # Default black (options don't have inherent color in provided HTML)
                 options_default_is_bold = False # Options container doesn't have font-weight in provided HTML

                 oc_style = options_container_div.get('style', '')
                 oc_parsed_size_px = parse_font_size_px(oc_style)
                 if oc_parsed_size_px:
                     options_default_size_pt = Pt(oc_parsed_size_px * 0.75)


                 option_y_offset_base = Inches(2.5) # Starting BASE Y position for the first option block
                 vertical_spacing_between_options = Inches(0.6) # Adjust this value for spacing


                 for j, option_element in enumerate(option_elements):
                      option_style = option_element.get('style', '')
                      translate_x_px, translate_y_px = parse_translate_px(option_style)

                      # Skip elements with 'noprint' class
                      if 'noprint' in option_element.get('class', []):
                           sys.stderr.write(f"  Slide {i+1}: Skipping 'noprint' option block {j+1}.\n")
                           continue

                      # Calculate base position for this option based on its index (for stacking)
                      base_left = Inches(0.7) # Options are indented slightly
                      base_top_calculated = option_y_offset_base.inches + (j * vertical_spacing_between_options.inches)
                      base_top = Inches(base_top_calculated)

                      width = Inches(8.5) # Adjust width for options
                      height = Inches(0.5) # Adjust height for options (can be adjusted after adding text)


                      # Calculate final position by adding the individual option's translate offset
                      final_left_inches = base_left.inches + (translate_x_px * PX_TO_INCHES)
                      final_top_inches = base_top.inches + (translate_y_px * PX_TO_INCHES)


                      # Extract formatted text runs from this option element, using options container defaults
                      # The extract_formatted_text_runs function will handle inline styles overriding these defaults.
                      option_text_runs = extract_formatted_text_runs(
                           option_element,
                           options_default_size_pt, # Pass options container default size
                           options_default_color_rgb, # Pass options container default color
                           options_default_is_bold # Pass options container default bold
                      )

                      if option_text_runs:
                           sys.stderr.write(f"  Slide {i+1}: Adding Option {chr(97+j)} shape.\n") # Debug print to stderr
                           txBox = slide.shapes.add_textbox(Inches(final_left_inches), Inches(final_top_inches), width, height)
                           tf = txBox.text_frame
                           tf.clear()
                           p = tf.add_paragraph()

                           for k, run_data in enumerate(option_text_runs):
                                run = p.add_run()
                                run.text = run_data['text']

                                font = run.font
                                # Apply formatting from the run_data (extracted from inline styles or block defaults)
                                font.size = run_data['font_size'] # Use the size determined by extract_formatted_text_runs
                                font.bold = run_data['is_bold'] # Use the bold determined by extract_formatted_text_runs
                                if run_data.get('color'):
                                    font.color.rgb = run_data['color'] # Use the color determined by extract_formatted_text_runs
                                # else: # Optional: Use the block default color if no inline color was found for this run
                                #      font.color.rgb = options_default_color_rgb # Use options container default

                                if run_data.get('is_superscript'):
                                    font.superscript = True
                                    # Optional: Reduce size for superscript if needed and not handled by inline style
                                    # if font.size is not None:
                                    #      font.size = Pt(font.size.pt * 0.7)


                 sys.stderr.write(f"Finished processing Options for Question Area {i+1}.\n") # Debug print to stderr


            sys.stderr.write(f"Finished processing Question Area {i+1}.\n") # Debug print to stderr


        # Save the presentation to a BytesIO object (in memory)
        sys.stderr.write("Saving presentation to BytesIO...\n") # Debug print to stderr
        binary_output = io.BytesIO()
        prs.save(binary_output)
        binary_output.seek(0) # Rewind to the beginning of the BytesIO object
        sys.stderr.write(f"Finished saving presentation to BytesIO. Bytes generated: {binary_output.getbuffer().nbytes}\n") # Debug print to stderr


        # Send the file back as a response
        sys.stderr.write("Sending PPTX file as response...\n") # Debug print to stderr
        return send_file(binary_output,
                         mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                         as_attachment=True,
                         download_name='generated_presentation.pptx')


    except Exception as e:
        # Catch any exceptions during the process and return an error response
        sys.stderr.write(f"Error during PPTX generation: {e}\n") # Debug print to stderr
        return jsonify({"error": str(e)}), 500


if __name__ == '__main__':
    # Run the Flask development server
    # In a production environment, you would use a more robust server like Gunicorn or uWSGI
    sys.stderr.write("Starting Flask development server...\n") # Debug print to stderr
    # app.run(debug=True) # debug=True prints errors to console, useful for development
    app.run(host='0.0.0.0', port=5000, debug=True) # Run on 0.0.0.0 to be accessible externally if needed, port 5000 is common