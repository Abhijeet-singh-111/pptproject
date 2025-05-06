from bs4 import BeautifulSoup, NavigableString, Tag # Import NavigableString and Tag
from pptx import Presentation
from pptx.util import Inches, Pt
import re
import html
from pptx.dml.color import RGBColor

# Your HTML content
html_string = """
<div class="rightarea order-2" id="mainPage" style="width: 100%; margin: auto;">
    <div class="QuestionArea">
      <div style="float:right;width:100%;" id="container">
        <div style="float:left;width:100%;color: red;border:solid 4px white;padding:4px;font-size:30px;font-weight:600;">
          <div class="alterable parent" style="transform: translate(100px, 300px); text-wrap: wrap;">
            <label class="alterable">Q. <span class="changeQNo" style="display:none;">1</span></label>&nbsp;
            <span id="fullElement" class="lng_english" style="display:none;color: blue !important;">
              11Find a quadratic polynomial whose sum and product of zeroes are 5, 7 respectively. </br></span>
            <span class="lng_hindi">एक द्विघात बहुपद ज्ञात कीजिए, जिसके शून्यांकों के योग तथा गुणनफल क्रमश: 5, 7 है।</span>
            <span class="noprint alterable" style="font-size:12px; transform: translate(-1056px, 44px);">6685954</span>
          </div>
        </div>

        <div class="alterable parent" style="width:90%;font-size:30px;font-weight:600;">
          <div class="alterable parent" style="display:inline-flex;margin-left: 5px;">
            <span class="alterable">(a)</span>&nbsp;
            <span class="lng_eng qoption_english optionPosition alterable" style="display:none;"><p>x11<sup>2</sup>&nbsp;+ 5x + 7</p></span>
            <span class="option_both_slash alterable" style="display:none;padding:0px 2px"> / </span>
            <span class="lng_hin qoption_hindi optionPosition alterable" style="display:inline;">
              <p>x<sup>2</sup>&nbsp;+ 5x + 7(hindi)</p>
            </span>
          </div>
            </br>
          <div class="alterable parent" style="display:inline-flex;margin-left: 5px;">
            <span class="alterable">(b)</span>&nbsp;
            <span class="lng_eng qoption_english optionPosition alterable" style="display:none;"><p>x11<sup>2</sup>&nbsp;+ 5x + 7</p></span>
            <span class="option_both_slash alterable" style="display:none;padding:0px 2px"> / </span>
            <span class="lng_hin qoption_hindi optionPosition alterable" style="display:inline;">
              <p>x<sup>2</sup>&nbsp;+ 5x + 7(hindi)</p>
            </span>
          </div>
            </br>
          <div class="alterable parent" style="display:inline-flex;margin-left: 5px;">
            <span class="alterable">(c)</span>&nbsp;
            <span class="lng_eng qoption_english optionPosition alterable" style="display:none;"><p>x11<sup>2</sup>&nbsp;+ 5x + 7</p></span>
            <span class="option_both_slash alterable" style="display:none;padding:0px 2px"> / </span>
            <span class="lng_hin qoption_hindi optionPosition alterable" style="display:inline;">
              <p>x<sup>2</sup>&nbsp;+ 5x + 7(hindi)</p>
            </span>
          </div>
        </div>
      </div>
    </div>

    <div class="QuestionArea">
      <div style="float:right;width:100%;" id="container">
        <div style="float:left;width:100%;color: red;border:solid 4px white;padding:4px;font-size:30px;font-weight:600;">
          <div class="alterable parent" style="transform: translate(0px, 0px); text-wrap: wrap;">
            <label class="alterable">Q. <span class="changeQNo" style="display:none;">1</span></label>&nbsp;
            <span id="fullElement" class="lng_english" style="display:none;color: blue !important;">
              22Find a quadratic polynomial whose sum and product of zeroes are 5, 7 respectively. </br></span>
            <span class="lng_hindi">एक द्विघात बहुपद ज्ञात कीजिए, जिसके शून्यांकों के योग तथा गुणनफल क्रमश: 5, 7 है।</span>
            <span class="noprint alterable" style="font-size:12px; transform: translate(-1056px, 44px);">6685954</span>
          </div>
        </div>

        <div class="alterable parent" style="width:90%;font-size:30px;font-weight:600;">
          <div class="alterable parent" style="display:inline-flex;margin-left: 5px;">
            <span class="alterable">(a)</span>&nbsp;
            <span class="lng_eng qoption_english optionPosition alterable" style="display:none;"><p>x11<sup>2</sup>&nbsp;+ 5x + 7</p></span>
            <span class="option_both_slash alterable" style="display:none;padding:0px 2px"> / </span>
            <span class="lng_hin qoption_hindi optionPosition alterable" style="display:inline;">
              <p>x<sup>2</sup>&nbsp;+ 5x + 7(hindi)</p>
            </span>
          </div>
            </br>
          <div class="alterable parent" style="display:inline-flex;margin-left: 5px;">
            <span class="alterable">(b)</span>&nbsp;
            <span class="lng_eng qoption_english optionPosition alterable" style="display:none;"><p>x11<sup>2</sup>&nbsp;+ 5x + 7</p></span>
            <span class="option_both_slash alterable" style="display:none;padding:0px 2px"> / </span>
            <span class="lng_hin qoption_hindi optionPosition alterable" style="display:inline;">
              <p>x<sup>2</sup>&nbsp;+ 5x + 7(hindi)</p>
            </span>
          </div>
            </br>
          <div class="alterable parent" style="display:inline-flex;margin-left: 5px;">
            <span class="alterable">(c)</span>&nbsp;
            <span class="lng_eng qoption_english optionPosition alterable" style="display:none;"><p>x11<sup>2</sup>&nbsp;+ 5x + 7</p></span>
            <span class="option_both_slash alterable" style="display:none;padding:0px 2px"> / </span>
            <span class="lng_hin qoption_hindi optionPosition alterable" style="display:inline;">
              <p>x<sup>2</sup>&nbsp;+ 5x + 7(hindi)</p>
            </span>
          </div>
        </div>
      </div>
    </div>
  </div>
"""

# Function to extract text and superscript info from a node, respecting display:none
def extract_text_runs_with_display_check(node):
    runs_data = []
    if node is None:
        return runs_data

    for content in node.contents:
        if isinstance(content, NavigableString): # This is a text node
            text = str(content)
            # Strip leading/trailing whitespace from the text node content
            cleaned_text = text.strip()
            # Replace internal sequences of whitespace (including newlines) with a single space
            cleaned_text = re.sub(r'\s+', ' ', cleaned_text)

            if cleaned_text: # Only add as a part if there's non-whitespace content
                 runs_data.append({'text': html.unescape(cleaned_text), 'is_superscript': False})

        elif isinstance(content, Tag): # This is an element node
            style = content.get('style', '')
            # Check if the element itself has display:none
            if 'display:none' in style.replace(" ", ""): # Replace space to handle 'display: none'
                 continue # Skip this element and its children if display is none

            # Handle specific tags or recurse for others
            if content.name == 'sup':
                sup_content_html = ''.join([str(c) for c in content.contents])
                sup_text_decoded = html.unescape(sup_content_html)
                if sup_text_decoded.strip():
                     runs_data.append({'text': sup_text_decoded, 'is_superscript': True})

            elif content.name == 'p':
                 # Recursively process content within <p>
                 p_content_runs = extract_text_runs_with_display_check(content) # Use new function
                 if p_content_runs:
                      runs_data.extend(p_content_runs)
                      # Option: Add a newline after <p> content if desired
                      # runs_data.append({'text': '\n'}) # Add a newline after <p> content

            # Add handling for other tags (b, i, br, etc.) and recurse for general tags
            # elif content.name == 'br':
            #     runs_data.append({'text': '\n'})
            else:
                # For other tags not explicitly handled, recurse into their children
                runs_data.extend(extract_text_runs_with_display_check(content)) # Use new function

    return runs_data


# Function to parse pixel font size from style string
def parse_font_size_px(style_string):
    font_size_match = re.search(r'font-size:\s*(\d+)px', style_string)
    if font_size_match:
        return int(font_size_match.group(1))
    return None

# Function to parse color from style string (basic)
def parse_color(style_string):
    color_match = re.search(r'color:\s*([^;]+)', style_string)
    if color_match:
        color_name_or_hex = color_match.group(1).strip()
        if color_name_or_hex.lower() == 'red':
             return RGBColor(0xFF, 0x00, 0x00)
        elif color_name_or_hex.startswith('#') and len(color_name_or_hex) == 7:
             try:
                 rgb_tuple = bytes.fromhex(color_name_or_hex[1:])
                 return RGBColor(rgb_tuple[0], rgb_tuple[1], rgb_tuple[2])
             except ValueError:
                 pass # Invalid hex
        # Add more color name mappings if needed
    return None # Default or if parsing failed

# Function to parse font weight from style string
def parse_font_weight(style_string):
    if 'font-weight:600' in style_string or 'font-weight:bold' in style_string:
        return True
    return False

# Function to parse translate(x, y) from transform style (assuming px)
def parse_translate_px(style_string):
    translate_match = re.search(r'transform:\s*translate\(\s*(-?\d+\.?\d*)px\s*,\s*(-?\d+\.?\d*)px\s*\)', style_string)
    if translate_match:
        try:
            translate_x = float(translate_match.group(1))
            translate_y = float(translate_match.group(2))
            return translate_x, translate_y
        except ValueError:
            pass # Could not convert to float
    return 0, 0 # Return 0,0 if not found or parsing failed


# --- Constants for Pixel to Inch Conversion (assuming 96 DPI) ---
PX_TO_INCHES = 0.1

# --- PPTX Generation ---
prs = Presentation()
blank_slide_layout = prs.slide_layouts[5]

# Parse the HTML
soup = BeautifulSoup(html_string, 'lxml')

# Find all QuestionArea divs
question_areas = soup.find_all('div', class_='QuestionArea')

if question_areas:
    print(f"Found {len(question_areas)} question areas.")
    for i, question_area in enumerate(question_areas):
        print(f"Processing Question Area {i+1}...")

        slide = prs.slides.add_slide(blank_slide_layout)

        # --- Extract Question Text ---
        # Find the main question div (with red color and font size)
        question_div = question_area.find('div', style=lambda value: value and 'color: red' in value and 'font-size:30px' in value)

        question_text_parts = []
        default_question_font_size_pt = Pt(20)
        question_color_rgb = RGBColor(0, 0, 0) # Default black
        question_is_bold = False

        # Initialize translate values before potential parsing
        question_translate_x_px, question_translate_y_px = 0, 0

        if question_div:
            # Extract style from the main question div for base styles
            question_div_style = question_div.get('style', '')
            parsed_size_px = parse_font_size_px(question_div_style)
            if parsed_size_px:
                 question_font_size_pt = Pt(parsed_size_px * 0.75)
            parsed_color = parse_color(question_div_style)
            if parsed_color:
                 question_color_rgb = parsed_color
            question_is_bold = parse_font_weight(question_div_style)

            # --- Find the specific alterable parent div for translate ---
            alterable_parent_div = question_div.find('div', class_='alterable parent')
            if alterable_parent_div:
                 alterable_parent_style = alterable_parent_div.get('style', '')
                 # Parse translate from THIS div's style
                 question_translate_x_px, question_translate_y_px = parse_translate_px(alterable_parent_style)
                 # --- Debugging: Print parsed translate for Question ---
                 print(f"  Question (Alterable Parent) translate: x={question_translate_x_px}px, y={question_translate_y_px}px")
                 # --- End Debugging ---
            else:
                 print("Warning: Could not find the 'alterable parent' div within the question div.")


             # Find the 'Q.' label (should be within the alterable parent div if found)
            q_label = None
            if alterable_parent_div:
                 q_label = alterable_parent_div.find('label', class_='alterable')
            elif question_div: # Fallback: try finding directly within question_div if alterable parent not found
                 q_label = question_div.find('label', class_='alterable')

            # --- Extract Q. text manually, respecting display:none ---
            q_text = ""
            if q_label:
                 # Use the new extraction function on the label content
                 q_label_content_runs = extract_text_runs_with_display_check(q_label)
                 # Join the text from the extracted runs
                 q_text = "".join(part['text'] for part in q_label_content_runs).strip()

                 # --- Debugging: Print extracted Q. label text ---
                 print(f"  Extracted Q. label text: '{q_text}'")
                 # --- End Debugging ---

            if q_text: # Only add the Q. part if extracted text is not empty
                 # Add Q. as a part, apply bold and size from question_div style
                 question_text_parts.append({'text': q_text + " ", 'is_bold': question_is_bold, 'font_size': question_font_size_pt})
            else:
                 print("Warning: Could not find or extract non-empty text from the 'Q.' label.")


             # Find the language spans (should be within the alterable parent div if found)
            english_span = None
            hindi_span = None
            if alterable_parent_div:
                 english_span = alterable_parent_div.find('span', class_='lng_english')
                 hindi_span = alterable_parent_div.find('span', class_='lng_hindi')
            elif question_div: # Fallback: try finding directly within question_div
                 english_span = question_div.find('span', class_='lng_english')
                 hindi_span = question_div.find('span', class_='lng_hindi')


            # Extract text from the visible language span using the new function
            main_question_content_runs = []
            selected_language_span = None # Track the selected span to get its style

            # Check display style to select the correct language span before processing
            if english_span and 'display:none' not in english_span.get('style', ''):
                 main_question_content_runs = extract_text_runs_with_display_check(english_span) # Use new function
                 selected_language_span = english_span
            elif hindi_span and ('display:inline' in hindi_span.get('style', '') or not hindi_span.get('style')):
                 main_question_content_runs = extract_text_runs_with_display_check(hindi_span) # Use new function
                 selected_language_span = hindi_span
            # Add else if for default language if neither has display:inline or style


            # Get style from the selected language span if available, otherwise use question_div style
            content_font_size_pt = question_font_size_pt
            content_color_rgb = question_color_rgb
            content_is_bold = question_is_bold # Inherit bold from question_div for content

            if selected_language_span:
                 lang_span_style = selected_language_span.get('style', '')
                 parsed_size_px = parse_font_size_px(lang_span_style)
                 if parsed_size_px:
                      content_font_size_pt = Pt(parsed_size_px * 0.75)
                 parsed_color = parse_color(lang_span_style)
                 if parsed_color:
                      content_color_rgb = parsed_color
                 # Note: Font weight on this span might override the parent div if present


            # --- Debugging: Print content runs for main question ---
            print(f"  Main Question Content Runs Debug:")
            for k, part in enumerate(main_question_content_runs):
                 print(f"    Part {k}: Text='{repr(part['text'])}', Superscript: {part.get('is_superscript', False)}")
            # --- End Debugging ---

            # Add the extracted runs from the selected language span
            for part in main_question_content_runs:
                 # Apply formatting derived from language span or question_div
                 part['font_size'] = content_font_size_pt
                 part['color'] = content_color_rgb
                 part['is_bold'] = part.get('is_bold', content_is_bold) # Allow part to override inherited bold
                 question_text_parts.append(part)


        # --- Add Question Shape ---
        if question_text_parts:
            # Define base position (adjust as needed)
            base_left = Inches(0.5)
            base_top = Inches(0.5)
            width = Inches(9)
            height = Inches(1.5) # Give more height

            # Apply translate offset (APPROXIMATION)
            offset_left_inches = question_translate_x_px * PX_TO_INCHES
            offset_top_inches = question_translate_y_px * PX_TO_INCHES

            # Fix: Ensure left and top are Inches objects by adding the offset to the base Inches values' magnitude
            left_inches_calculated = base_left.inches + offset_left_inches
            top_inches_calculated = base_top.inches + offset_top_inches

            left = Inches(left_inches_calculated)
            top = Inches(top_inches_calculated)


            # --- Debugging: Print final calculated position for Question ---
            print(f"  Question Final Position: left={left.inches:.2f}in, top={top.inches:.2f}in")
            # --- End Debugging ---


            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.clear()

            p = tf.add_paragraph()

            # Add text runs with formatting based on the prepared parts
            for part in question_text_parts:
                 run = p.add_run()
                 run.text = part['text']
                 font = run.font
                 font.size = part.get('font_size', default_question_font_size_pt)
                 font.bold = part.get('is_bold', False)
                 if part.get('color'):
                      font.color.rgb = part['color']
                 # Superscript is handled by extract_text_runs and would be applied here if the part had it.


        # --- Extract and Add Options ---
        options_div = question_area.find('div', style=lambda value: value and 'width:90%' in value and 'font-size:30px' in value)

        if options_div:
            option_elements = options_div.find_all('div', style=lambda value: value and 'display:inline-flex' in value)
            option_y_offset_base = Inches(2.5) # Starting BASE Y position (adjust)

            # Default styles from options_div
            options_div_style = options_div.get('style', '')
            options_font_size_pt = Pt(20) # Initialized here
            parsed_size_px = parse_font_size_px(options_div_style)
            if parsed_size_px:
                 options_font_size_pt = Pt(parsed_size_px * 0.75)
            # Color from options_div is less relevant as color is on question_div, but could be parsed here too.
            # Font-weight from options_div is also less relevant as letter is bolded and content follows span style.


            for j, option_element in enumerate(option_elements):
                # Parse translate from the option_element div style
                option_translate_x_px, option_translate_y_px = parse_translate_px(option_element.get('style', ''))
                 # --- Debugging: Print parsed translate for Option ---
                print(f"  Option {chr(97+j)} translate: x={option_translate_x_px}px, y={option_translate_y_px}px")
                 # --- End Debugging ---


                option_letter_span = option_element.find('span', class_='alterable')
                option_letter = option_letter_span.get_text(strip=True) if option_letter_span else ''

                # Manually find the language spans within this specific option element
                english_option_span = None
                hindi_span = None
                for child_span in option_element.find_all('span', recursive=False): # Only look at direct children
                     if isinstance(child_span, Tag) and child_span.name == 'span': # Ensure it's a tag and a span
                          classAttribute = child_span.get('class', [])
                          if 'lng_eng' in classAttribute:
                               english_option_span = child_span
                          if 'lng_hin' in classAttribute:
                               hindi_span = child_span
                     # Ignore other direct children like the follow-btn div


                # Extract text from the visible language span using the new function
                option_content_parts = []
                selected_option_language_span = None

                # Check display style to select the correct language span before processing
                if english_option_span and 'display:none' not in english_option_span.get('style', ''):
                     option_content_parts.extend(extract_text_runs_with_display_check(english_option_span)) # Use new function
                     selected_option_language_span = english_option_span
                elif hindi_span and ('display:inline' in hindi_span.get('style', '') or not hindi_span.get('style')):
                     option_content_parts.extend(extract_text_runs_with_display_check(hindi_span)) # Use new function
                     selected_option_language_span = hindi_span
                # Add else if for default language if neither has display:inline or style


                # Get style from the selected language span if available, otherwise use options_div style
                option_content_font_size_pt = options_font_size_pt # Use the correctly initialized variable
                # option_content_color_rgb = ... # Add color logic if needed from span
                # option_content_is_bold = ... # Add bold logic if needed from span

                if selected_option_language_span:
                     lang_span_style = selected_option_language_span.get('style', '')
                     parsed_size_px = parse_font_size_px(lang_span_style)
                     if parsed_size_px:
                          option_content_font_size_pt = Pt(parsed_size_px * 0.75)
                     # Parse color or bold from span style if they override parent div


                # --- Add Option Shape ---
                if option_letter or option_content_parts:
                    # Define base position (adjust as needed)
                    base_left = Inches(0.7)
                    # Fix: Calculate base_top for option ensuring it's an Inches object
                    base_top_inches_calculated = option_y_offset_base.inches + (j * 0.6)
                    base_top_option = Inches(base_top_inches_calculated)


                    width = Inches(8.5) # Original width
                    # Suggestion: Temporarily increase width to test if text wraps
                    # width = Inches(12) # <-- Uncomment this line to test if a wider box fixes the single line issue
                    height = Inches(0.5)


                    # Apply translate offset (APPROXIMATION)
                    offset_left_inches = option_translate_x_px * PX_TO_INCHES
                    offset_top_inches = option_translate_y_px * PX_TO_INCHES

                    # Fix: Ensure left and top are Inches objects by adding the offset to the base Inches values' magnitude
                    left_inches_calculated = base_left.inches + offset_left_inches
                    top_inches_calculated = base_top_option.inches + offset_top_inches # Use the correctly calculated base_top for option

                    left = Inches(left_inches_calculated)
                    top = Inches(top_inches_calculated)

                    # --- Debugging: Print final calculated position for Option ---
                    print(f"  Option {chr(97+j)} Final Position: left={left.inches:.2f}in, top={top.inches:.2f}in")
                    # --- End Debugging ---


                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.clear()

                    p = tf.add_paragraph()

                    # Add the option letter run
                    if option_letter:
                        run = p.add_run()
                        run.text = option_letter + ". "
                        font = run.font
                        font.bold = True
                        font.size = options_font_size_pt


                    # Add the option text content runs
                    for part in option_content_parts:
                         run = p.add_run()
                         run.text = part['text']
                         font = run.font

                         font.size = option_content_font_size_pt

                         if part.get('is_superscript'):
                              font.superscript = True
                              font.size = Pt(option_content_font_size_pt.pt * 0.7)

                         # Apply bold/color here if extracted


                    # Vertical spacing is controlled by the base_top calculation


# Save the presentation
output_filename = "presentation_python_translated.pptx"
prs.save(output_filename)

print(f"\nPresentation saved to {output_filename}")